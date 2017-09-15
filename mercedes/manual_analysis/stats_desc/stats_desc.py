import os
import logging
# import seaborn as sns
import matplotlib as mpl
import matplotlib.pyplot as plt
from statsmodels.graphics.mosaicplot import mosaic

import pandas as pd
import datetime
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Matplotlib parameters ------------------------------------------------------------------------------------------------
# No spines
mpl.rcParams['axes.facecolor'] = 'white'
mpl.rcParams['axes.spines.left'] = False
mpl.rcParams['axes.spines.right'] = False
mpl.rcParams['axes.spines.top'] = False

# No tick
mpl.rcParams['xtick.major.size'] = 0
mpl.rcParams['ytick.major.size'] = 0

mpl.rcParams['axes.edgecolor'] = 'grey'

# mpl.rcParams['figure.dpi'] = 100
mpl.rcParams['figure.figsize'] = 4.10, 3.46

mpl.rcParams['savefig.dpi'] = 300
mpl.rcParams['savefig.bbox'] = 'tight'
mpl.rcParams['savefig.transparent'] = True

mpl.rcParams['boxplot.patchartist'] = True
mpl.rcParams['boxplot.boxprops.color'] = '#7F7F7F'
mpl.rcParams['boxplot.whiskerprops.color'] = '#218DAA'
mpl.rcParams['boxplot.medianprops.color'] = '#7F7F7F'

# Manage outliers in boxplot
mpl.rcParams['boxplot.flierprops.marker'] = 'o'
mpl.rcParams['boxplot.flierprops.markerfacecolor'] = '#218DAA'
mpl.rcParams['boxplot.flierprops.markeredgecolor'] = '#7F7F7F'
mpl.rcParams['boxplot.flierprops.markersize'] = 4

# ----------------------------------------------------------------------------------------------------------------------

class Logger(object):
    ''' Logger object to use in functions
    '''
    def __init__(self):
        pass

    def create_logger(self):
        '''
            Initialise logger object with file and console handlers.
        '''
        logger = logging.getLogger(__name__)
        logger.setLevel(logging.INFO)

        # # create a file handler
        # fh = logging.FileHandler('hello.log')
        # fh.setLevel(logging.INFO)

        # create console handler (with a higher log level)
        ch = logging.StreamHandler()
        ch.setLevel(logging.INFO)

        # create a logging format
        formatter_f = logging.Formatter(
            '%(asctime)s - %(name)s - %(levelname)s - %(message)s')
        formatter_c = logging.Formatter('%(levelname)s - %(message)s')
        # fh.setFormatter(formatter_f)
        ch.setFormatter(formatter_c)

        # add the handlers to the logger
        # logger.addHandler(fh)
        logger.addHandler(ch)
        return logger


class PowerPointFileManager(object):
    ''' Class to manage PowerPoint documents
    '''
    def __init__(self):
        '''
            Class constructor
            Includes variables to manage paths and creates folders if they do not exist
        '''
        self.PPT_TEMPLATE = "1. inputs/input_stats_desc.pptx"
        # Slide layout: 0: Slicer; 1: 2 columns with table placeholder on RHS
        self.ppt_output = "../stats_desc.pptx"

        self.CM_TO_INCH = 0.393701
        pass

    def duplicate_ppt_template(self):
        '''
            Duplicate PowerPoint template and create empty output document

            Parameters
            ----------
            None

            Returns
            -------
            None
        '''
        prs = Presentation(self.PPT_TEMPLATE)
        prs.save(self.ppt_output)

        pass

    def add_slicer(self, slicer_number, section_name):
        '''
            Create a slicer slide with slicer_number and section_name.

            Parameters
            ----------
            slicer_number: string
                Section number to print on the left-hand side.
            section_name: string
                Section name to print on the right-hand side.

            Returns
            -------
            None
        '''
        prs = Presentation(self.ppt_output)
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        placeholder_number = slide.placeholders[18]
        placeholder_number.text = slicer_number

        placeholder_section_name = slide.placeholders[17]
        placeholder_section_name.text = section_name

        prs.save(self.ppt_output)
        pass

    def add_slide_with_empty_table(self, var_type):
        '''
            Create a slide with an empty stats table on the right-hand side.

            Parameters
            ----------
            var_type: string {'categorical', 'boolean', 'int', 'float', 'numeric'} or numpy type

            Returns
            -------
            None
        '''
        prs = Presentation(self.ppt_output)
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # Placeholder 18 is a Table Placeholder
        placeholder = slide.placeholders[18]
        graphic_frame = placeholder.insert_table(rows=10, cols=2)
        table_inserted = graphic_frame.table

        # Transparent background for all cells of the table
        # Format content : size = 12, font color = grey ----------------------------------------------------------------
        for i in range(10):
            row = table_inserted.rows[i]
            row.height = Inches(0.30)

            for j in range(2):
                cell = table_inserted.cell(i, j)
                cell.fill.background()
                cell.margin_left = 0

                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.font.bold = False
                paragraph.font.color.rgb = RGBColor(0, 0, 0)

                if j == 1:
                    paragraph.alignment = PP_ALIGN.RIGHT
        # --------------------------------------------------------------------------------------------------------------

        # Texts in the cells of the table ------------------------------------------------------------------------------
        table_inserted.cell(2, 0).text = 'Taux NaN'
        table_inserted.cell(3, 0).text = 'Nombre de valeurs uniques'

        if var_type in ['object', 'categorical', 'boolean', 'bool']:
            table_inserted.cell(6, 0).text = 'Mode'
        elif var_type in ['int16', 'int32', 'int64', 'float16', 'float32', 'float64', 'int', 'float', 'numeric']:
            table_inserted.cell(5, 0).text = 'Min'
            table_inserted.cell(6, 0).text = 'Mediane'
            table_inserted.cell(7, 0).text = 'Max'

        table_inserted.cell(9, 0).text = 'Exemples'
        # --------------------------------------------------------------------------------------------------------------

        prs.save(self.ppt_output)
        pass

    def add_empty_slide_2_columns(self):
        '''
            Create a 2-column slide with subtitle placeholders.

            Parameters
            ----------
            None

            Returns
            -------
            None
        '''
        prs = Presentation(self.ppt_output)
        slide = prs.slides.add_slide(prs.slide_layouts[2])

        # TODO: add placeholders and 2 arguments placeholder_left, placeholder_right

        prs.save(self.ppt_output)
        pass

    def add_slide_title(self, var):
        '''
            Add title on slide

            Parameters
            ----------
            var: string
                Name of variable to get info from

            Returns
            -------
            None
        '''
        prs = Presentation(self.ppt_output)
        slide = prs.slides[len(prs.slides) - 1]

        title = slide.placeholders[0]
        title.text = 'Statistiques descriptives – Variable ' + var

        prs.save(self.ppt_output)
        pass

    def add_graph_lhs(self, output_directory, filename):
        '''
            Add filename on left-hand side on last slide of presentation_name

            Parameters
            ----------
            filename: string
                File to add

            Returns
            -------
            None
        '''
        prs = Presentation(self.ppt_output)
        slide = prs.slides[len(prs.slides) - 1]

        picture = slide.shapes.add_picture(output_directory + filename,
                                           left=Inches(1.3 * self.CM_TO_INCH),
                                           top=Inches(4.19 * self.CM_TO_INCH),
                                           width=Inches(10.44 * self.CM_TO_INCH))

        prs.save(self.ppt_output)
        pass

    def add_graph_rhs(self, output_directory, filename):
        '''
            Add filename on right-hand side on last slide of presentation_name

            Parameters
            ----------
            filename: string
                File to add

            Returns
            -------
            None
        '''
        prs = Presentation(self.ppt_output)
        slide = prs.slides[len(prs.slides) - 1]

        picture = slide.shapes.add_picture(output_directory + filename,
                                           left=Inches(13.67 * self.CM_TO_INCH),
                                           top=Inches(4.19 * self.CM_TO_INCH),
                                           width=Inches(10.44 * self.CM_TO_INCH))

        prs.save(self.ppt_output)
        pass

    def add_info_table_rhs(self, df, var, var_type):
        '''
            Add data analyses on right-hand side on last slide of presentation_name:
                * NaN rate
                * Number of unique values
                * Min, Median, Max
                * Value examples

            Parameters
            ----------
            df: dataframe
                Dataframe containing var at least
            var: string
                Name of variable to get info from
            var_type: string {'categorical', 'boolean', 'int', 'float', 'numeric'} or numpy type

            Returns
            -------
            None
        '''
        prs = Presentation(self.ppt_output)
        slide = prs.slides[len(prs.slides) - 1]

        table = slide.placeholders[18].table
        table.cell(2, 1).text = str(round(100 * df[var].isnull().sum() / df.shape[0], 2)) + ' %'
        table.cell(3, 1).text = str(len(df[var].unique()))

        if var_type in ['object', 'categorical', 'boolean', 'bool']:
            table.cell(6, 1).text = str(df[var].mode()[0])
        elif var_type in ['int16', 'int32', 'int64', 'float16', 'float32', 'float64', 'int', 'float', 'numeric']:
            table.cell(5, 1).text = str(np.min(df[var]))
            table.cell(6, 1).text = str(np.median(df[var]))
            table.cell(7, 1).text = str(np.max(df[var]))

        table.cell(9, 1).text = str(df[var].unique()[0:5].tolist())

        prs.save(self.ppt_output)
        pass


class StatsDesc(object):
    '''
        Main class for package stats_desc
    '''

    def __init__(self):
        '''
            Class constructor
        '''
        self.output_directory = '3. outputs/'
        self.color = '#218DAA'
        self.color_secondary = '#1F497D'

        if not os.path.exists(self.output_directory):
            os.mkdir(self.output_directory)

        self.logger = Logger()
        self.ppt_doc = PowerPointFileManager()

        pass

    def save_eda(self, df, list_variables, target, target_type,
                 autodetect_types=True, type_variables=[], filename_prefix='',
                 bool_convert_percent=True,
                 bins=20, x_axis_range=None):
        '''
            Plot distribution of categorical variable column_name.

            Parameters
            ----------
            df: dataframe
                Dataframe containing list_variables at least.
            type_variables: list, {'categorical', 'boolean', 'int', 'float', 'numeric'}
            # TODO

            Returns
            -------
            None
        '''
        logger = self.logger.create_logger()

        # Types of variables -------------------------------------------------------------------------------------------
        logger.info('In progress: detection of variable types')

        if autodetect_types:
            types = df[list_variables].dtypes.values
        else:
            if len(list_variables) != len(type_variables):
                logger.info('Lengths of list_variables and type_variables not the same - Types are autodetected')
                types = df[list_variables].dtypes.values
            else:
                types = type_variables

        logger.info('Done: variable types detected')
        # --------------------------------------------------------------------------------------------------------------

        # Plots to save as png figures ---------------------------------------------------------------------------------
        logger.info('In progress: Monovariate analysis - saving plots as png figures')

        for index, col in enumerate(list_variables):
            if str(types[index]) in ['object', 'categorical', 'boolean', 'bool']:
                self.plot_univariate_categorical(df=df,
                                                 column_name=col,
                                                 filename_prefix=filename_prefix,
                                                 bool_convert_percent=bool_convert_percent)
            elif str(types[index]) in ['int16', 'int32', 'int64', 'float16', 'float32', 'float64', 'int', 'float', 'numeric']:
                self.plot_univariate_numeric(df=df,
                                             column_name=col,
                                             filename_prefix=filename_prefix,
                                             bins=bins,
                                             x_axis_range=x_axis_range)

        logger.info('Done: Monovariate analysis - plots saved as png figures')
        # --------------------------------------------------------------------------------------------------------------

        # Creation of document -----------------------------------------------------------------------------------------
        logger.info('In progress: Monovariate analysis - creation of PowerPoint document')

        self.ppt_doc.duplicate_ppt_template()

        self.ppt_doc.add_slicer(slicer_number='1', section_name='Analyse univariée')

        for index, col in enumerate(list_variables):
            self.ppt_doc.add_slide_with_empty_table(var_type=str(types[index]))
            self.ppt_doc.add_slide_title(var=col)
            self.ppt_doc.add_graph_lhs(output_directory=self.output_directory,
                                       filename=filename_prefix + 'distrib_' + col + '.png')
            self.ppt_doc.add_info_table_rhs(df=df, var=col, var_type=str(types[index]))

        logger.info('Done: Monovariate analysis - PowerPoint document created')
        # --------------------------------------------------------------------------------------------------------------

        # Plots to save as png figures ---------------------------------------------------------------------------------
        logger.info('In progress: Bivariate analysis - saving plots as png figures')

        if target_type not in ['numeric', 'categorical']:
            logger.error("target_type not 'numeric' or 'categorical'")
        else:
            for col in [col for col in list_variables if col != target]:
                index = list_variables.index(col)

                if target_type == 'numeric':
                    if str(types[index]) in ['object', 'categorical', 'boolean', 'bool']:
                        self.plot_bivariate_x_categorical_y_numeric(df=df, x_name=col, target_name=target,
                                                                    filename_prefix=filename_prefix)
                    elif str(types[index]) in ['int16', 'int32', 'int64', 'float16', 'float32', 'float64', \
                                               'int', 'float', 'numeric']:
                        self.plot_bivariate_x_numeric_y_numeric(df=df, x_name=col, target_name=target,
                                                                filename_prefix=filename_prefix)
                elif target_type == 'categorical':
                    if str(types[index]) in ['object', 'categorical', 'boolean', 'bool']:
                        self.plot_bivariate_x_categorical_y_categorical(df=df, x_name=col, target_name=target,
                                                                        filename_prefix=filename_prefix)
                    # TODO: to be coded
                    # elif str(types[index]) in ['int16', 'int32', 'int64', 'float16', 'float32', 'float64', \
                    #                            'int', 'float', 'numeric']:

        logger.info('Done: Bivariate analysis - plots saved as png figures')
        # --------------------------------------------------------------------------------------------------------------

        # Creation of document -----------------------------------------------------------------------------------------
        logger.info('In progress: Bivariate analysis - creation of PowerPoint document')

        self.ppt_doc.add_slicer(slicer_number='2', section_name='Analyse bivariée')

        if target_type not in ['numeric', 'categorical']:
            logger.error("target_type not 'numeric' or 'categorical'")
        else:
            for col in [col for col in list_variables if col != target]:
                index = list_variables.index(col)

                self.ppt_doc.add_empty_slide_2_columns()
                self.ppt_doc.add_slide_title(var='{} en fonction de {}'.format(target, col))

                if target_type == 'numeric':
                    if str(types[index]) in ['object', 'categorical', 'boolean', 'bool']:
                        self.ppt_doc.add_graph_lhs(output_directory=self.output_directory,
                                                   filename='{prefix}bivariate_boxplot_with_outlier_{y}_{x}.png'.\
                                                    format(prefix=filename_prefix, y=target, x=col))
                        self.ppt_doc.add_graph_rhs(output_directory=self.output_directory,
                                                   filename='{prefix}bivariate_boxplot_no_outlier_{y}_{x}.png'.\
                                                    format(prefix=filename_prefix, y=target, x=col))
                    elif str(types[index]) in ['int16', 'int32', 'int64', 'float16', 'float32', 'float64', \
                                               'int', 'float', 'numeric']:
                        self.ppt_doc.add_graph_lhs(output_directory=self.output_directory,
                                                   filename='{prefix}bivariate_bubbles_{y}_{x}.png'.\
                                                    format(prefix=filename_prefix, y=target, x=col))
                        self.ppt_doc.add_graph_rhs(output_directory=self.output_directory,
                                                   filename='{prefix}bivariate_scatter_{y}_{x}.png'.\
                                                    format(prefix=filename_prefix, y=target, x=col))
                elif target_type == 'categorical':
                    if str(types[index]) in ['object', 'categorical', 'boolean', 'bool']:
                        self.ppt_doc.add_graph_lhs(output_directory=self.output_directory,
                                                   filename='{prefix}bivariate_stacked_{y}_{x}.png'.\
                                                    format(prefix=filename_prefix, y=target, x=col))
                        self.ppt_doc.add_graph_rhs(output_directory=self.output_directory,
                                                   filename='{prefix}bivariate_mekko_{y}_{x}.png'.\
                                                    format(prefix=filename_prefix, y=target, x=col))

        logger.info('Done: Bivariate analysis - PowerPoint document created')
        # --------------------------------------------------------------------------------------------------------------

        logger.handlers.pop(0)
        pass

    def plot_univariate_categorical(self, df, column_name, filename_prefix='', bool_convert_percent=True):
        '''
            Plot distribution of categorical variable column_name.

            Parameters
            ----------
            df: dataframe
                Dataframe containing column_name at least.
            column_name: string
                Name of column to be plotted.
            filename_prefix: string
                Prefix added to filename.
            bool_convert_percent: boolean, default True
                Convert frequency values into percentage.

            Returns
            -------
            None
        '''
        df_to_plot = df[column_name].value_counts()
        df_to_plot.sort_values(ascending=True, inplace=True)

        fig, ax = plt.subplots()
        ax = df_to_plot.plot(kind='barh', alpha=0.6, color=self.color)

        # ax.set_title('Distribution de la variable ' + column_name)
        plt.title('Distribution de la variable {}'.format(column_name), y=1.08)

        ax.set_xlabel('Fréquence')
        ax.set_xticks([])

        ax.spines['left'].set_visible(True)
        ax.spines['bottom'].set_visible(False)

        # Labels to print on bar plot ----------------------------------------------------------------------------------
        rects = ax.patches
        labels = [x for x in df_to_plot.values]

        if bool_convert_percent:
            # labels = [str(int(100 * label / sum(labels))) + ' %' for label in labels]
            labels = ['{:.0%}'.format(label / sum(labels)) for label in labels]
            # TODO: change scale of y axis to 100 %

        for index, (rect, label) in enumerate(zip(rects, labels)):
            width = rect.get_width()
            ax.text(width + 20,
                    index,
                    label,
                    # ha='center',
                    va='center',
                    size=9)
        # --------------------------------------------------------------------------------------------------------------
        plt.savefig(self.output_directory + filename_prefix + 'distrib_' + column_name + '.png')
        plt.close(fig)

        pass

    # TODO: add bool_convert_percent?
    def plot_univariate_numeric(self, df, column_name, filename_prefix='', bins=20, x_axis_range=None):
        '''
            Plot distribution of numeric variable column_name.

            Parameters
            ----------
            df: dataframe
                Dataframe containing column_name at least.
            column_name: string
                Name of column to be plotted.
            filename_prefix: string
                Prefix added to filename.
            bins: integer or list, default 20
                If an integer, divide the counts in the specified number of bins.
                If a sequence of values, the values of the lower bound of the bins to be used.
            x_axis_range: list of 2 values, default None
                List [xmin, xmax] to resize x axis.

            Returns
            -------
            None
        '''
        df_to_plot = df[column_name]

        fig, ax = plt.subplots()

        ax = df_to_plot.plot(kind='hist', alpha=0.6, color=self.color, bins=bins, xlim=x_axis_range)

        # ax.set_title('Distribution de la variable ' + column_name)
        plt.title('Distribution de la variable {}'.format(column_name), y=1.08)
        ax.set_ylabel('Fréquence')

        plt.savefig(self.output_directory + filename_prefix + 'distrib_' + column_name + '.png')
        plt.close(fig)

        pass

    def plot_bivariate_x_categorical_y_categorical(self, df, x_name, target_name, filename_prefix=''):
        '''
            Plot bivariate analysis : y = f(x) where both x and y are categorical.
            This functions generates two graphs: a mekko chart and a stacked bar chart

            Parameters
            ----------
            df: dataframe
                Dataframe containing x_name and target_name at least.
            x_name: string
                Name of column that is on x axis.
            target_name: string
                Name of column containing target to predict.
            filename_prefix: string
                Prefix added to filename.

            Returns
            -------
            None
        '''
        df_to_plot = self._build_dataset_for_x_cat_y_cat(df=df, x_name=x_name, target_name=target_name)

        # Plot 1: Mekko graph ------------------------------------------------------------------------------------------
        df_to_plot['label'] = df_to_plot['count_percent'].apply(int).apply(str) + ' %' + '\n' \
                                + '(' + df_to_plot['count'].apply(str) + ')'

        props = lambda index: {'color': self.color if index[1] == str(df[target_name].value_counts().index[1]) \
                                                    else self.color_secondary,
                               'alpha': 0.7}
        labels = lambda k: df_to_plot.loc[k, 'label']
        plot = mosaic(data=df_to_plot['count'], gap=0.02,
                      title='Distribution de {target} en fonction de {var}'.format(target=target_name, var=x_name),
                      properties=props, labelizer=labels)

        plt.savefig(self.output_directory + filename_prefix + 'bivariate_mekko_' + target_name + '_' + x_name + '.png')
        plt.close()
        # --------------------------------------------------------------------------------------------------------------

        # Plot 2: Stacked bar chart ------------------------------------------------------------------------------------
        df_to_plot.reset_index(inplace=True)
        df_to_plot2 = df_to_plot.pivot(index=x_name, columns=target_name, values='count')
        df_to_plot2['total'] = df_to_plot2.sum(axis=1)
        df_to_plot2.sort_values(by='total', ascending=False, inplace=True)

        fig, ax = plt.subplots()

        bar_width = 0.75
        bar_position = [i+1 for i in range(df_to_plot2.shape[0])]
        tick_position = [i+(bar_width/2) for i in bar_position]

        ax.bar(bar_position,
               df_to_plot2[str(df[target_name].value_counts().index[0])],
               width=bar_width,
               label=str(df[target_name].value_counts().index[0]),
               alpha=0.7,
               color=self.color_secondary,
               edgecolor='#7F7F7F')

        ax.bar(bar_position,
               df_to_plot2[str(df[target_name].value_counts().index[1])],
               width=bar_width,
               bottom=df_to_plot2[str(df[target_name].value_counts().index[0])],
               label=str(df[target_name].value_counts().index[1]),
               alpha=0.7,
               color=self.color,
               edgecolor='#7F7F7F')

        rects = ax.patches
        labels_numbers = df_to_plot2[str(df[target_name].value_counts().index[0])].tolist() + \
                            df_to_plot2[str(df[target_name].value_counts().index[1])].tolist()

        labels_percent = (df_to_plot2[str(df[target_name].value_counts().index[0])]/df_to_plot2['total']).tolist() + \
                    (df_to_plot2[str(df[target_name].value_counts().index[1])]/df_to_plot2['total']).tolist()

        labels = ['{}\n({:.0%})'.format(number, percent) for number, percent in zip(labels_numbers, labels_percent)]

        for rect, label in zip(rects, labels):
            coord = rect.get_xy()
            height = rect.get_height()
            ax.text(x=coord[0] + bar_width/2, y=coord[1] + height/2, s=label,
                    size=9, va='center', ha='center', color='w')

        for index in range(len(df_to_plot2.index)):
            ax.text(x=bar_position[index] + bar_width/2,
                    y=df_to_plot2['total'].values[index],
                    s=df_to_plot2['total'].values[index],
                    size=10, ha='center', va='bottom')

        plt.xticks(tick_position, df_to_plot2.index.values)
        plt.yticks([])
        plt.xlim([min(tick_position)-bar_width, max(tick_position)+bar_width])

        plt.title('Distribution de {target} en fonction de {var}'.format(target=target_name, var=x_name), y=1.08)

        handles, labels = ax.get_legend_handles_labels()
        ax.legend(handles[::-1], labels[::-1], loc='upper right', frameon=False)

        plt.savefig('{}{}bivariate_stacked_{}_{}.png'.format(self.output_directory,
                                                             filename_prefix, target_name, x_name))
        plt.close(fig)
        # --------------------------------------------------------------------------------------------------------------
        pass

    def _build_dataset_for_x_cat_y_cat(self, df, x_name, target_name):
        '''
            Build the dataset for bivariate plot when both x and y are categorical.

            Parameters
            ----------
            df: dataframe
                Dataframe containing x_name and target_name at least.
            x_name: string
                Name of column that is on x axis.
            target_name: string
                Name of column containing target to predict.

            Returns
            -------
            df_with_stats: dataframe
            Index: [x_name, target_name]
            Columns:
                * count: number of observations for a given combination of x and y
                * count_percent: percentage of target distribution for every category of x
                * count_by_x: number of observations of x
        '''
        df_to_plot = df.copy(deep=True)
        df_to_plot[target_name] = df_to_plot[target_name].apply(str)
        df_to_plot['count'] = 1
        df_to_plot = pd.DataFrame(df_to_plot.groupby([x_name, target_name])['count'].count())

        df_count_percent = df_to_plot.groupby(level=0).apply(lambda x: round(100 * x / float(x.sum())))
        df_count_percent.columns = ['count_percent']

        df_count_by_x = df_to_plot.groupby(level=0).apply(sum)
        df_count_by_x.columns = ['count_by_x']

        df_merged = pd.merge(left=df_to_plot, left_index=True, right=df_count_by_x, right_index=True, how='left')
        df_sorted = df_merged.sort_values(by=['count_by_x', 'count'], ascending=False)
        df_with_stats = pd.merge(left=df_sorted, left_index=True, right=df_count_percent, right_index=True, how='left')

        return df_with_stats

    def plot_bivariate_x_numeric_y_numeric(self, df, x_name, target_name, filename_prefix=''):
        '''
            Plot bivariate analysis : y = f(x) where both x and y are numeric.
            This functions generates two graphs: a bubbles plot and a scatter plot.

            Parameters
            ----------
            df: dataframe
                Dataframe containing x_name and target_name at least.
            x_name: string
                Name of column that is on x axis.
            target_name: string
                Name of column containing target to predict.
            filename_prefix: string
                Prefix added to filename.

            Returns
            -------
            None
        '''
        col = df.columns[0]

        df_grouped = pd.DataFrame(df.groupby([x_name, target_name])[col].count())
        df_grouped.columns = ['size']
        df_grouped.reset_index(inplace=True)

        # Plot 1: Bubbles plot -----------------------------------------------------------------------------------------
        fig, ax = plt.subplots()
        ax = df_grouped.plot(kind='scatter', x=x_name, y=target_name,
                             s=df_grouped['size'], alpha=0.6, color=self.color)
        ax.spines['left'].set_visible(True)
        plt.title('Distribution de {target} en fonction de {x}'.format(target=target_name, x=x_name),
                  y=1.08)
        plt.savefig('{}{}bivariate_bubbles_{}_{}.png'.format(self.output_directory,
                                                             filename_prefix, target_name, x_name))
        plt.close(fig)
        # --------------------------------------------------------------------------------------------------------------

        # Plot 2: Scatter plot -----------------------------------------------------------------------------------------
        fig, ax = plt.subplots()
        ax = df_grouped.plot(kind='scatter', x=x_name, y=target_name,
                             alpha=0.6, color=self.color)
        ax.spines['left'].set_visible(True)
        plt.title('Distribution de {target} en fonction de {x}'.format(target=target_name, x=x_name),
                  y=1.08)
        plt.savefig('{}{}bivariate_scatter_{}_{}.png'.format(self.output_directory,
                                                             filename_prefix, target_name, x_name))
        plt.close(fig)
        # --------------------------------------------------------------------------------------------------------------
        pass

    def plot_bivariate_x_categorical_y_numeric(self, df, x_name, target_name, filename_prefix=''):
        df_to_plot = df.copy(deep=True)
        # Such a hard list so that x axis is ordered by decreasing order
        data = [df_to_plot[df_to_plot[x_name]==level][target_name].values \
                for level in df_to_plot[x_name].value_counts().index.tolist()]

        # Boxplot 1: No outliers ---------------------------------------------------------------------------------------
        fig, ax = plt.subplots()
        box = plt.boxplot(data, labels=df_to_plot[x_name].value_counts().index.tolist(), showfliers=False)

        # To color the boxplot
        colors = [self.color for i in range(len(data))]
        for patch, color in zip(box['boxes'], colors):
            patch.set_facecolor(color)
            patch.set_alpha(0.6)

        plt.title('Distribution de {target} en fonction de {x}'.format(target=target_name, x=x_name),
                  y=1.08)
        plt.savefig('{}{}bivariate_boxplot_no_outlier_{}_{}.png'.format(self.output_directory,
                                                                        filename_prefix, target_name, x_name))
        plt.close(fig)
        # --------------------------------------------------------------------------------------------------------------

        # Boxplot 1: With outliers -------------------------------------------------------------------------------------
        fig, ax = plt.subplots()
        box = plt.boxplot(data, labels=df_to_plot[x_name].value_counts().index.tolist(), showfliers=True)

        # To color the boxplot
        colors = [self.color for i in range(len(data))]
        for patch, color in zip(box['boxes'], colors):
            patch.set_facecolor(color)
            patch.set_alpha(0.6)

        plt.title('Distribution de {target} en fonction de {x}'.format(target=target_name, x=x_name),
                  y=1.08)
        plt.savefig('{}{}bivariate_boxplot_with_outlier_{}_{}.png'.format(self.output_directory,
                                                                          filename_prefix, target_name, x_name))
        plt.close(fig)
        # --------------------------------------------------------------------------------------------------------------

        pass



if __name__ == '__main__':
    print('ok')
